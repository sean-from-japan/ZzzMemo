#!/usr/bin/env python3
"""
docs/generate_docs.py — qcatch 使い方ガイド (PPTX) を生成する

使い方: python docs/generate_docs.py
出力:   docs/qcatch_guide.pptx

必要: pip install python-pptx
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("python-pptx が必要です。インストール:  pip install python-pptx")
    raise SystemExit(1)

import io, sys
from pathlib import Path

if sys.platform == "win32":
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

OUT = Path(__file__).parent / "qcatch_guide.pptx"

# ── カラーパレット ────────────────────────────────────────────────────────────
CYN = RGBColor(0x00, 0xB4, 0xD8)
CYN_DK = RGBColor(0x00, 0x77, 0xB6)
DARK = RGBColor(0x1E, 0x20, 0x2B)
DIM = RGBColor(0x6C, 0x75, 0x7D)
WHT = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xE5, 0x38, 0x35)
GRN = RGBColor(0x28, 0xA7, 0x45)
YLW = RGBColor(0xFD, 0xA8, 0x0B)
BG = RGBColor(0xF0, 0xFB, 0xFF)
CODE = RGBColor(0x2B, 0x2D, 0x42)
ORG = RGBColor(0xFF, 0x7F, 0x00)

W, H = Inches(13.33), Inches(7.5)


# ── ヘルパー ─────────────────────────────────────────────────────────────────
def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, fill):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    return shape


def tb(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    sz=18,
    bold=False,
    color=DARK,
    italic=False,
    align=PP_ALIGN.LEFT,
    font="Segoe UI",
    wrap=True,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def code_block(slide, text, x, y, w, h, sz=13):
    rect(slide, x, y, w, h, CODE)
    tb(
        slide,
        text,
        x + Inches(0.18),
        y + Inches(0.12),
        w - Inches(0.36),
        h - Inches(0.24),
        sz=sz,
        color=CYN,
        font="Courier New",
        wrap=False,
    )


def header_bar(slide, title, emoji=""):
    rect(slide, 0, 0, W, Inches(0.85), CYN_DK)
    label = f"{emoji}  {title}" if emoji else title
    tb(
        slide,
        label,
        Inches(0.4),
        Inches(0.1),
        W - Inches(0.8),
        Inches(0.68),
        sz=26,
        bold=True,
        color=WHT,
    )


def bullets(slide, items, x, y, w, h, sz=17, color=DARK):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(7)
        run = p.add_run()
        run.text = f"•   {item}"
        run.font.name = "Segoe UI"
        run.font.size = Pt(sz)
        run.font.color.rgb = color


# ── スライド定義 ─────────────────────────────────────────────────────────────


def slide_title(prs):
    s = new_slide(prs)
    rect(s, 0, 0, W, H, DARK)
    rect(s, 0, H - Inches(0.7), W, Inches(0.7), CYN_DK)
    tb(
        s,
        "⚡  qcatch",
        0,
        Inches(1.3),
        W,
        Inches(1.5),
        sz=64,
        bold=True,
        color=CYN,
        align=PP_ALIGN.CENTER,
    )
    tb(
        s,
        "爆速タスクキャッチ & AI 自動整理",
        0,
        Inches(3.0),
        W,
        Inches(0.8),
        sz=28,
        color=WHT,
        align=PP_ALIGN.CENTER,
    )
    tb(
        s,
        "CLI タスク管理ツール  —  Windows 11  /  Python",
        0,
        Inches(3.85),
        W,
        Inches(0.6),
        sz=19,
        color=DIM,
        align=PP_ALIGN.CENTER,
    )
    tb(
        s,
        "python-pptx で自動生成 — docs/generate_docs.py",
        0,
        H - Inches(0.58),
        W,
        Inches(0.4),
        sz=11,
        color=RGBColor(0x55, 0x66, 0x77),
        align=PP_ALIGN.CENTER,
        italic=True,
    )


def slide_concept(prs):
    s = new_slide(prs)
    header_bar(s, "コンセプト", "💡")

    tb(
        s,
        "課題",
        Inches(0.6),
        Inches(1.0),
        Inches(5.6),
        Inches(0.45),
        sz=17,
        bold=True,
        color=RED,
    )
    bullets(
        s,
        [
            "タスクが頭に浮かんでもメモアプリを開くのが面倒",
            "後で「あれ何だっけ？」となる",
            "TODO リストが散在して整理できない",
        ],
        Inches(0.6),
        Inches(1.55),
        Inches(5.6),
        Inches(2.0),
        sz=16,
    )

    tb(
        s,
        "解決策",
        Inches(0.6),
        Inches(3.65),
        Inches(5.6),
        Inches(0.45),
        sz=17,
        bold=True,
        color=GRN,
    )
    bullets(
        s,
        [
            "思いついた瞬間にトースト通知から記録（0.5秒）",
            "時間があるときに AI が自動分類して整理",
        ],
        Inches(0.6),
        Inches(4.2),
        Inches(5.6),
        Inches(1.4),
        sz=16,
    )

    # フロー図（右カラム）
    rect(s, Inches(7.0), Inches(1.0), Inches(6.0), Inches(6.0), BG)
    tb(
        s,
        "思いついた！",
        Inches(7.2),
        Inches(1.2),
        Inches(5.6),
        Inches(0.5),
        sz=18,
        bold=True,
        color=DARK,
        align=PP_ALIGN.CENTER,
    )

    rect(s, Inches(7.9), Inches(1.85), Inches(4.2), Inches(0.72), CYN_DK)
    tb(
        s,
        "Windows Search → qcatch",
        Inches(7.9),
        Inches(1.85),
        Inches(4.2),
        Inches(0.72),
        sz=15,
        bold=True,
        color=WHT,
        align=PP_ALIGN.CENTER,
    )

    tb(
        s,
        "↓  Quick Add ウィンドウが出る",
        Inches(7.2),
        Inches(2.68),
        Inches(5.6),
        Inches(0.4),
        sz=13,
        color=DIM,
        align=PP_ALIGN.CENTER,
    )

    rect(s, Inches(7.9), Inches(3.15), Inches(4.2), Inches(0.62), DIM)
    tb(
        s,
        "inbox.txt  に自動追記",
        Inches(7.9),
        Inches(3.15),
        Inches(4.2),
        Inches(0.62),
        sz=14,
        bold=True,
        color=WHT,
        align=PP_ALIGN.CENTER,
        font="Courier New",
    )

    tb(
        s,
        "↓  時間があるとき",
        Inches(7.2),
        Inches(3.88),
        Inches(5.6),
        Inches(0.4),
        sz=13,
        color=DIM,
        align=PP_ALIGN.CENTER,
    )

    rect(s, Inches(7.9), Inches(4.35), Inches(4.2), Inches(0.72), ORG)
    tb(
        s,
        "python qcatch.py sort",
        Inches(7.9),
        Inches(4.35),
        Inches(4.2),
        Inches(0.72),
        sz=14,
        bold=True,
        color=WHT,
        align=PP_ALIGN.CENTER,
        font="Courier New",
    )

    tb(
        s,
        "↓  AI が自動分類",
        Inches(7.2),
        Inches(5.18),
        Inches(5.6),
        Inches(0.4),
        sz=13,
        color=DIM,
        align=PP_ALIGN.CENTER,
    )

    rect(s, Inches(7.9), Inches(5.65), Inches(4.2), Inches(0.62), GRN)
    tb(
        s,
        "sorted_tasks.md  に保存",
        Inches(7.9),
        Inches(5.65),
        Inches(4.2),
        Inches(0.62),
        sz=14,
        bold=True,
        color=WHT,
        align=PP_ALIGN.CENTER,
        font="Courier New",
    )


def slide_toast(prs):
    s = new_slide(prs)
    header_bar(s, "最速入力  —  toast コマンド（Quick Add ウィンドウ）", "⚡")

    tb(
        s,
        "起動フロー",
        Inches(0.6),
        Inches(1.0),
        Inches(12.0),
        Inches(0.45),
        sz=17,
        bold=True,
        color=CYN_DK,
    )

    steps = [
        ("Windows キー", "「qcatch」と入力 → Enter"),
        ("Quick Add ウィンドウ", "タスクを入力 → Enter で追加"),
        ("ループ継続 or 終了", "次を入力 or 空Enter / Esc で閉じる"),
    ]
    for i, (title, desc) in enumerate(steps):
        x = Inches(0.6) + i * Inches(4.1)
        rect(s, x, Inches(1.55), Inches(3.9), Inches(1.4), BG)
        rect(s, x, Inches(1.55), Inches(0.5), Inches(1.4), [CYN_DK, ORG, GRN][i])
        tb(
            s,
            str(i + 1),
            x,
            Inches(1.55),
            Inches(0.5),
            Inches(1.4),
            sz=22,
            bold=True,
            color=WHT,
            align=PP_ALIGN.CENTER,
        )
        tb(
            s,
            title,
            x + Inches(0.6),
            Inches(1.65),
            Inches(3.2),
            Inches(0.5),
            sz=15,
            bold=True,
            color=DARK,
        )
        tb(
            s,
            desc,
            x + Inches(0.6),
            Inches(2.2),
            Inches(3.2),
            Inches(0.65),
            sz=13,
            color=DIM,
        )

    tb(
        s,
        "UI の特徴",
        Inches(0.6),
        Inches(3.2),
        Inches(5.8),
        Inches(0.45),
        sz=17,
        bold=True,
        color=CYN_DK,
    )
    bullets(
        s,
        [
            "最前面表示・起動直後にフォーカス自動取得",
            "Shift + Enter で改行（複数行 → 複数タスクに分割保存）",
            "追加後フィールドをクリア → そのまま次のタスクを入力可能",
            "追加件数をステータスラベルに 3 秒表示",
            "sv-ttk（Windows 11 風テーマ）対応",
        ],
        Inches(0.6),
        Inches(3.75),
        Inches(5.8),
        Inches(3.0),
        sz=14,
    )

    tb(
        s,
        "コマンド例",
        Inches(7.0),
        Inches(3.2),
        Inches(5.9),
        Inches(0.45),
        sz=17,
        bold=True,
        color=CYN_DK,
    )
    code_block(
        s,
        "# Windows Search から起動\n"
        "qcatch  →  Enter\n\n"
        "# ターミナルから直接\n"
        "python qcatch.py toast\n"
        'python qcatch.py add "タスク"  # API なし・即追記',
        Inches(7.0),
        Inches(3.75),
        Inches(5.9),
        Inches(2.4),
        sz=13,
    )

    tb(
        s,
        "tkinter 標準ライブラリで実装 — 追加依存なし・--noconsole ビルド対応",
        Inches(0.6),
        Inches(6.85),
        Inches(12.0),
        Inches(0.38),
        sz=13,
        color=DIM,
        italic=True,
    )


def slide_dashboard(prs):
    s = new_slide(prs)
    header_bar(s, "タスク閲覧・管理  —  dashboard コマンド", "📊")

    # 左カラム: Inbox タブ
    tb(
        s,
        "Inbox タブ（未分類）",
        Inches(0.6),
        Inches(1.0),
        Inches(5.8),
        Inches(0.45),
        sz=17,
        bold=True,
        color=CYN_DK,
    )
    bullets(
        s,
        [
            "inbox.txt の全タスクをリスト表示",
            "タスクを選択 → Delete キーで削除（確認ダイアログあり）",
            "件数をタブタイトルにリアルタイム表示",
        ],
        Inches(0.6),
        Inches(1.55),
        Inches(5.8),
        Inches(2.0),
        sz=15,
    )

    rect(s, Inches(0.6), Inches(3.65), Inches(5.8), Inches(2.85), BG)
    tb(
        s,
        "Inbox (3 件)",
        Inches(0.8),
        Inches(3.8),
        Inches(5.4),
        Inches(0.4),
        sz=13,
        bold=True,
        color=CYN_DK,
    )
    for j, t in enumerate(
        [
            "[2026-03-20 09:00] 企画書を月曜までに提出",
            "[2026-03-20 09:05] 牛乳・卵を買う",
            "[2026-03-20 09:10] Python 非同期処理を学ぶ",
        ]
    ):
        rect(
            s,
            Inches(0.8),
            Inches(4.25) + j * Inches(0.6),
            Inches(5.4),
            Inches(0.52),
            BG if j % 2 == 0 else WHT,
        )
        tb(
            s,
            t,
            Inches(0.9),
            Inches(4.3) + j * Inches(0.6),
            Inches(5.2),
            Inches(0.42),
            sz=11,
            color=DARK,
            font="Courier New",
        )

    # 右カラム: 分類済みタブ
    tb(
        s,
        "分類済みタブ（sorted_tasks.md）",
        Inches(7.0),
        Inches(1.0),
        Inches(5.9),
        Inches(0.45),
        sz=17,
        bold=True,
        color=CYN_DK,
    )
    bullets(
        s,
        [
            "sorted_tasks.md をスクロール表示（リードオンリー）",
            "カテゴリ（仕事 / 買い物 / 学習 …）をそのまま確認",
            "編集が必要な場合はテキストエディタで直接 .md を開く",
        ],
        Inches(7.0),
        Inches(1.55),
        Inches(5.9),
        Inches(2.0),
        sz=15,
    )

    code_block(
        s,
        "# 分類済みタスク表示例\n"
        "# タスク整理 (2026-03-20 09:30)\n\n"
        "## 仕事\n"
        "- [09:00] 企画書を月曜までに提出\n\n"
        "## 買い物\n"
        "- [09:05] 牛乳・卵を買う\n\n"
        "## 学習\n"
        "- [09:10] Python 非同期処理を学ぶ",
        Inches(7.0),
        Inches(3.65),
        Inches(5.9),
        Inches(2.85),
        sz=12,
    )

    tb(
        s,
        "起動: python qcatch.py dashboard",
        Inches(0.6),
        Inches(6.85),
        Inches(12.0),
        Inches(0.38),
        sz=13,
        bold=True,
        color=CYN_DK,
        font="Courier New",
    )


def slide_sort(prs):
    s = new_slide(prs)
    header_bar(s, "AI 自動整理  —  sort コマンド", "🤖")

    # 左カラム: コマンド
    tb(
        s,
        "コマンド",
        Inches(0.6),
        Inches(1.0),
        Inches(6.2),
        Inches(0.45),
        sz=17,
        bold=True,
        color=CYN_DK,
    )
    code_block(
        s,
        "# Gemini 2.5 Flash（推奨・無料）\n"
        "python qcatch.py sort\n\n"
        "# Ollama（ローカル・完全無料・オフライン）\n"
        "python qcatch.py sort --local\n\n"
        "# API なし（プロンプト書き出し）\n"
        "python qcatch.py sort --export",
        Inches(0.6),
        Inches(1.55),
        Inches(6.2),
        Inches(3.3),
        sz=13,
    )

    tb(
        s,
        "出力: data/sorted_tasks.md",
        Inches(0.6),
        Inches(5.0),
        Inches(6.2),
        Inches(0.45),
        sz=15,
        bold=True,
        color=CYN_DK,
    )
    code_block(
        s,
        "# タスク整理 (2026-03-20 09:00)\n\n"
        "## 仕事\n- [09:00] 企画書を月曜までに提出\n\n"
        "## 買い物\n- [08:55] 牛乳、卵\n\n"
        "## 学習\n- [08:58] Pythonの非同期処理",
        Inches(0.6),
        Inches(5.55),
        Inches(6.2),
        Inches(1.6),
        sz=12,
    )

    # 右カラム: バックエンド比較
    rect(s, Inches(7.3), Inches(1.0), Inches(5.7), Inches(6.0), BG)
    tb(
        s,
        "バックエンド優先順位",
        Inches(7.5),
        Inches(1.15),
        Inches(5.3),
        Inches(0.45),
        sz=16,
        bold=True,
        color=CYN_DK,
    )

    backends = [
        (GRN, "① --local", "Ollama（phi4 等）", "完全ローカル・オフライン・無料"),
        (CYN, "② Gemini", "GEMINI_API_KEY 設定済み", "2.5 Flash・無料枠・最速"),
        (
            YLW,
            "③ Anthropic",
            "ANTHROPIC_API_KEY 設定済み",
            "Claude Haiku・有料フォールバック",
        ),
        (ORG, "④ --export", "API 不使用", "プロンプトをファイル出力"),
    ]
    for i, (col, flag, cond, desc) in enumerate(backends):
        ry = Inches(1.75) + i * Inches(1.15)
        rect(s, Inches(7.5), ry, Inches(0.45), Inches(0.85), col)
        tb(
            s,
            flag,
            Inches(8.05),
            ry + Inches(0.05),
            Inches(4.6),
            Inches(0.45),
            sz=15,
            bold=True,
            color=DARK,
            font="Courier New",
        )
        tb(
            s,
            cond,
            Inches(8.05),
            ry + Inches(0.5),
            Inches(4.6),
            Inches(0.38),
            sz=13,
            color=DIM,
        )

    tb(
        s,
        "Pydantic 構造化出力（JSON Schema）で分類精度を向上\ntemperature=0.1 でカテゴリのブレを最小化",
        Inches(7.5),
        Inches(6.3),
        Inches(5.3),
        Inches(0.6),
        sz=13,
        color=DIM,
        italic=True,
    )


def slide_setup(prs):
    s = new_slide(prs)
    header_bar(s, "セットアップ", "⚙")

    for i, (badge_col, step, code, note) in enumerate(
        [
            (
                CYN_DK,
                "① パッケージインストール",
                "pip install -r requirements.txt",
                None,
            ),
            (
                CYN_DK,
                "② .exe をビルドして Windows Search に登録",
                "Set-ExecutionPolicy -ExecutionPolicy RemoteSigned -Scope CurrentUser\n.\\build.ps1",
                "qcatch.exe（全機能統合・39MB）を生成してスタートメニューに登録",
            ),
            (
                GRN,
                "③ Gemini API キーを取得して設定（無料・クレカ不要）",
                '[Environment]::SetEnvironmentVariable("GEMINI_API_KEY", "AIzaSy...", "User")',
                "取得先: aistudio.google.com/app/apikey  ─  PowerShell 再起動で反映",
            ),
        ]
    ):
        base_y = Inches(1.0) + i * Inches(1.9)
        rect(
            s, Inches(0.5), base_y + Inches(0.05), Inches(0.35), Inches(0.45), badge_col
        )
        tb(
            s,
            step,
            Inches(1.0),
            base_y,
            Inches(12.0),
            Inches(0.5),
            sz=16,
            bold=True,
            color=DARK,
        )
        code_block(
            s,
            code,
            Inches(0.5),
            base_y + Inches(0.55),
            Inches(12.1),
            Inches(0.82 if "\n" in code else Inches(0.5)),
            sz=12,
        )
        if note:
            tb(
                s,
                note,
                Inches(0.5),
                base_y + Inches(1.5),
                Inches(12.1),
                Inches(0.38),
                sz=13,
                color=DIM,
                italic=True,
            )

    tb(
        s,
        "Ollama（ローカル LLM）を使う場合は ollama.com からアプリをインストール後: ollama pull phi4",
        Inches(0.5),
        Inches(6.85),
        Inches(12.3),
        Inches(0.38),
        sz=13,
        color=DIM,
        italic=True,
    )


def slide_architecture(prs):
    s = new_slide(prs)
    header_bar(s, "アーキテクチャ", "🏗")

    code_block(
        s,
        "ZzzMemo/\n"
        "├── qcatch.py             # メインスクリプト（全コマンド）\n"
        "├── qcatch.exe            # PyInstaller でビルド（全機能統合）\n"
        "├── qcatch_config.json    # 設定ファイル（sort_backend 等）\n"
        "├── build.ps1             # .exe ビルド + Windows Search 登録\n"
        "├── requirements.txt\n"
        "│\n"
        "├── data/                 # 実行時データ\n"
        "│   ├── inbox.txt         # 未整理タスク（タイムスタンプ付き）\n"
        "│   ├── sorted_tasks.md   # AI 分類済みタスク\n"
        "│   ├── archive.txt       # sort 済みの inbox バックアップ\n"
        "│   └── qcatch.log        # GUI 操作ログ（自動生成）\n"
        "│\n"
        "├── docs/\n"
        "│   ├── generate_docs.py  # この PPTX を生成するスクリプト\n"
        "│   ├── windows_setup.md  # Windows セットアップ手順\n"
        "│   └── qcatch_guide.pptx\n"
        "│\n"
        "└── prompts/              # Gemini Q&A プロンプトと回答\n"
        "    ├── gemini_ui_improvement.md\n"
        "    ├── gemini_ui_followup.md\n"
        "    └── archive/          # 実装済みの過去プロンプト",
        Inches(0.5),
        Inches(1.0),
        Inches(8.3),
        Inches(6.1),
        sz=12,
    )

    tb(
        s,
        "設計上のポイント",
        Inches(9.1),
        Inches(1.0),
        Inches(4.0),
        Inches(0.45),
        sz=16,
        bold=True,
        color=CYN_DK,
    )
    bullets(
        s,
        [
            "add / toast / dashboard は API 通信なし（待ち時間ゼロ）",
            "sort のみ外部 API または Ollama を使用",
            "GUI は tkinter で統一  ─  --noconsole ビルド・sv-ttk テーマ対応",
            "データはテキストファイルのまま維持  ─  他ツール連携・可読性のため",
            "sort 実行は inbox.txt をクリアする  ─  事前に dashboard で確認推奨",
        ],
        Inches(9.1),
        Inches(1.6),
        Inches(4.0),
        Inches(4.5),
        sz=14,
    )


def slide_commands(prs):
    s = new_slide(prs)
    header_bar(s, "コマンドリファレンス", "📋")

    rows = [
        ("toast", "", "Quick Add ウィンドウを起動（ループ入力・最前面表示）"),
        ("dashboard", "", "ダッシュボード（Inbox 管理 + 分類済みタスク閲覧）"),
        ("add", '"タスクのテキスト"', "即追記（API通信なし・待ち時間ゼロ）"),
        ("list", "", "inbox の一覧表示"),
        ("prompt", "", "ターミナル対話入力"),
        ("sort", "", "Gemini 2.5 Flash で分類（GEMINI_API_KEY 必要）"),
        ("sort --local", "", "Ollama（ローカル LLM）で分類"),
        ("sort --export", "", "API 不使用・プロンプトを data/sort_prompt.txt に出力"),
    ]

    hy = Inches(1.05)
    row_h = 0.62
    rect(s, Inches(0.4), hy, Inches(12.5), Inches(0.5), CYN_DK)
    for lbl, cx, cw in [
        ("コマンド", Inches(0.6), Inches(3.5)),
        ("引数", Inches(4.3), Inches(2.5)),
        ("説明", Inches(7.0), Inches(5.8)),
    ]:
        tb(s, lbl, cx, hy + Inches(0.07), cw, Inches(0.38), sz=15, bold=True, color=WHT)

    for i, (cmd, arg, desc) in enumerate(rows):
        ry = hy + Inches(0.5) + i * Inches(row_h)
        rect(
            s,
            Inches(0.4),
            ry,
            Inches(12.5),
            Inches(row_h - 0.02),
            BG if i % 2 == 0 else WHT,
        )
        tb(
            s,
            cmd,
            Inches(0.6),
            ry + Inches(0.13),
            Inches(3.5),
            Inches(0.45),
            sz=14,
            bold=True,
            color=CYN_DK,
            font="Courier New",
        )
        tb(
            s,
            arg,
            Inches(4.3),
            ry + Inches(0.13),
            Inches(2.5),
            Inches(0.45),
            sz=13,
            color=DIM,
            font="Courier New",
        )
        tb(
            s,
            desc,
            Inches(7.0),
            ry + Inches(0.13),
            Inches(5.8),
            Inches(0.45),
            sz=14,
            color=DARK,
        )

    tb(
        s,
        '起動方法: Windows Search → 「qcatch」Enter  /  PowerShell → qcatch add "..."',
        Inches(0.4),
        Inches(6.65),
        Inches(12.5),
        Inches(0.45),
        sz=14,
        color=DIM,
        italic=True,
        align=PP_ALIGN.CENTER,
    )


# ── エントリポイント ─────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slide_title(prs)
    slide_concept(prs)
    slide_toast(prs)
    slide_dashboard(prs)
    slide_sort(prs)
    slide_setup(prs)
    slide_architecture(prs)
    slide_commands(prs)

    prs.save(OUT)
    print(f"✅ 生成完了: {OUT}  ({len(prs.slides)} スライド)")


if __name__ == "__main__":
    main()
